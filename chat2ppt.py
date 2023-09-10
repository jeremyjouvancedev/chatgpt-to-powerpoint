import json
import argparse

# Env variable
from dotenv import load_dotenv

# PowerPoint Handling
from pptx import Presentation
from pptx.util import Inches

# Langchain Handling
from langchain.chat_models import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain.chains.openai_functions import (
    create_openai_fn_chain,
    create_structured_output_chain,
)
from langchain.prompts import ChatPromptTemplate, HumanMessagePromptTemplate
from jsonSchema import json_schema

load_dotenv()  # take environment variables from .env.

def generate_langchain_chain(nb_of_slides=10):
    chat_model = ChatOpenAI(verbose=True)

    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "You are a world class Powerpoint generator."),
            ("human", "Generate a " + str(nb_of_slides) + " slide Powerpoint presentation for the topic: {input}"),
            ("human", "Tip: Make sure to answer in the correct format"),
        ]
    )

    chain = create_structured_output_chain(json_schema, chat_model, prompt, verbose=True)

    return chain

def layout_title_and_content(prs, slide_data):
    slide_layout = prs.slide_layouts[1]  # 'Title and Content' layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    slide.shapes.title.text = slide_data["title"]
    
    # Add content to the slide
    content_shape = slide.placeholders[1]
    for point in slide_data["content"]:
        p = content_shape.text_frame.add_paragraph()
        p.text = point    


def generate_slides(nb_of_slides, topic, output, load_json = False):
    response = None

    if not load_json:
        chain = generate_langchain_chain(nb_of_slides)
        response = chain.run(topic)

        with open(output+'.json', 'w') as json_file:
            json.dump(response, json_file)
    else:
        with open(output + '.json', 'r') as json_file:
            response = json.load(json_file)

    # Create a new presentation object
    prs = Presentation()

    slides_data = response['slides']

    for idx, slide_data in enumerate(slides_data):
        print(f"Slide {idx+1}/{len(slides_data)}")

        if slide_data['layout'] == 'Title Slide':
            slide_layout = prs.slide_layouts[0]  # 'Title and Content' layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            slide.shapes.title.text = slide_data["title"]


        elif slide_data['layout'] == 'Title and Content':
            layout_title_and_content(prs, slide_data)

        elif slide_data['layout'] == 'Two Content':
            if len(slide_data.get('content', [])) > 0 and len(slide_data.get('left_content', [])) == 0 and len(slide_data.get('right_content', [])) == 0:
                layout_title_and_content(prs, slide_data)        
            else:
                slide_layout = prs.slide_layouts[2]
                slide = prs.slides.add_slide(slide_layout)

                # Set slide title
                slide.shapes.title.text = slide_data["title"]

                # Add content to the left placeholder
                left_content_shape = slide.placeholders[1]
                for point in slide_data["left_content"]:
                    p = left_content_shape.text_frame.add_paragraph()
                    p.text = point
                
                # Add content to the right placeholder
                right_content_shape = slide.placeholders[2]
                for point in slide_data["right_content"]:
                    p = right_content_shape.text_frame.add_paragraph()
                    p.text = point
        
        elif slide_data['layout'] == 'Comparison':
            slide_layout = prs.slide_layouts[3]  # Assuming 'Comparison' layout is at index 3
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            slide.shapes.title.text = slide_data["title"]
            
            # Add content to the left placeholder
            left_content_shape = slide.placeholders[1]
            for point in slide_data["left_content"]:
                p = left_content_shape.text_frame.add_paragraph()
                p.text = point
            
            # Add content to the right placeholder
            right_content_shape = slide.placeholders[2]
            for point in slide_data["right_content"]:
                p = right_content_shape.text_frame.add_paragraph()
                p.text = point

        elif slide_data['layout'] == 'Content with Caption':
            slide_layout = prs.slide_layouts[4]  # Assuming 'Content with Caption' layout is at index 4
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            slide.shapes.title.text = slide_data["title"]
            
            # Add content
            content_shape = slide.placeholders[2]
            for point in slide_data["content"]:
                p = content_shape.text_frame.add_paragraph()
                p.text = point
            
            # Add caption
            caption_shape = slide.placeholders[4]
            caption_shape.text = slide_data["caption"]

        else:    
            print(f'Layout >{slide_data["layout"]}< not implemented')

    # Save the presentation
    prs.save(output)

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentation using ChatGPT.')
    parser.add_argument('topic', type=str, help='Topic for the PowerPoint presentation.')
    parser.add_argument('--slides', type=int, default=10, help='Number of slides in the presentation. Default is 10.')
    parser.add_argument('--output', type=str, default='presentation.pptx', help='Output PowerPoint file name. Default is presentation.pptx.')
    args = parser.parse_args()

    generate_slides(args.slides, args.topic, args.output)