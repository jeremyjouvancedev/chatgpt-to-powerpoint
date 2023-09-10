from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
from langchain.chat_models import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain.chains.openai_functions import (
    create_openai_fn_chain,
    create_structured_output_chain,
)
from langchain.prompts import ChatPromptTemplate, HumanMessagePromptTemplate
from jsonSchema import json_schema

load_dotenv()  # take environment variables from .env.

NB_OF_SLIDES = 10

chat_model = ChatOpenAI(verbose=True)


prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are a world class Powerpoint generator."),
        ("human", "Generate a " + NB_OF_SLIDES + " slide Powerpoint presentation for the topic: {input}"),
        ("human", "Tip: Make sure to answer in the correct format"),
    ]
)


chain = create_structured_output_chain(json_schema, chat_model, prompt, verbose=True)
response = chain.run("chatgpt")

print(response)


# Create a new presentation object
prs = Presentation()

slides_data = response['slides']

for slide_data in slides_data:
    slide_layout = prs.slide_layouts[1]  # 'Title and Content' layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    slide.shapes.title.text = slide_data["title"]
    
    # Add content to the slide
    content_shape = slide.placeholders[1]
    for point in slide_data["content"]:
        p = content_shape.text_frame.add_paragraph()
        p.text = point

# Save the presentation
prs.save('presentation.pptx')